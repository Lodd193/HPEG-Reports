#!/usr/bin/env python3
"""
HPEG Report Generator - Script 2 of 2
======================================
Generates NHS-branded PowerPoint presentations for HPEG reporting.

This script:
1. Loads processed data from hpeg_data_processor.py
2. Creates professional charts using matplotlib
3. Generates PowerPoint presentations (one per HPEG)
4. Applies NHS branding and formatting

Author: Patient Relations Quality & Performance Manager
Version: 1.0
Date: December 2024
"""

import sys
import pickle
from pathlib import Path
from datetime import datetime
import tempfile

# Fix Windows console encoding for Unicode characters
if sys.stdout.encoding != 'utf-8':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
        sys.stderr.reconfigure(encoding='utf-8')
    except:
        pass

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# USER CONFIGURATION - UPDATE THESE PATHS FOR YOUR WORK ENVIRONMENT
# ============================================================================

PROCESSED_DATA_PATH = r"C:/Users/lod19/HPEGs/processed_data.pkl"
OUTPUT_FOLDER = r"C:/Users/lod19/HPEGs/outputs"

# ============================================================================
# NHS BRANDING CONFIGURATION
# ============================================================================

NHS_COLORS_RGB = {
    'nhs_blue': (0, 94, 184),
    'nhs_dark_blue': (0, 48, 135),
    'nhs_bright_blue': (0, 114, 206),
    'nhs_light_blue': (65, 182, 230),
    'nhs_aqua_blue': (0, 169, 206),
    'nhs_black': (35, 31, 32),
    'nhs_dark_grey': (66, 85, 99),
    'nhs_mid_grey': (118, 134, 146),
    'nhs_pale_grey': (232, 237, 238),
    'white': (255, 255, 255),
    'nhs_dark_green': (0, 103, 71),
    'nhs_green': (0, 150, 57),
    'nhs_light_green': (120, 190, 32),
    'nhs_aqua_green': (0, 164, 153),
    'nhs_purple': (51, 0, 114),
    'dark_pink': (124, 40, 85),
    'nhs_pink': (174, 37, 115),
    'nhs_dark_red': (138, 21, 56),
    'emergency_red': (218, 41, 28),
    'nhs_orange': (237, 139, 0),
    'nhs_warm_yellow': (255, 184, 28),
    'nhs_yellow': (250, 225, 0)
}

# Matplotlib color codes (hex)
NHS_COLORS_HEX = {
    'nhs_blue': '#005EB8',
    'nhs_dark_blue': '#003087',
    'nhs_bright_blue': '#0072CE',
    'nhs_light_blue': '#41B6E6',
    'nhs_aqua_blue': '#00A9CE',
    'nhs_dark_grey': '#425563',
    'nhs_mid_grey': '#768692',
    'nhs_pale_grey': '#E8EDEE',
    'nhs_black': '#231F20',
    'white': '#FFFFFF',
    'nhs_dark_green': '#006747',
    'nhs_green': '#009639',
    'nhs_light_green': '#78BE20',
    'nhs_aqua_green': '#00A499',
    'nhs_purple': '#330072',
    'dark_pink': '#7C2855',
    'nhs_pink': '#AE2573',
    'nhs_dark_red': '#8A1538',
    'emergency_red': '#DA291C',
    'nhs_orange': '#ED8B00',
    'nhs_warm_yellow': '#FFB81C',
    'nhs_yellow': '#FAE100'
}

# ============================================================================
# POWERPOINT HELPER FUNCTIONS
# ============================================================================

def create_presentation():
    """Create a blank widescreen PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    return prs

def add_title_bar(slide, hpeg_name, slide_title):
    """
    Add NHS Blue title bar to slide.
    Format: "HPEG NAME | Slide Title"
    """
    # Full-width blue rectangle
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=0,
        width=Inches(13.333),
        height=Inches(0.8)
    )

    # Style the rectangle
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_blue'])
    title_bar.line.fill.background()

    # Add text
    text_frame = title_bar.text_frame
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.3)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = f"{hpeg_name.upper()} | {slide_title}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'

    return title_bar

def add_text_box(slide, text, left, top, width, height, font_size=11, bold=False,
                 color_rgb=NHS_COLORS_RGB['nhs_black'], alignment=PP_ALIGN.LEFT):
    """Add a text box to the slide."""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color_rgb)
    p.font.name = 'Arial'
    p.alignment = alignment

    return textbox

def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add an image to the slide."""
    if width and height:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top),
                                width=Inches(width), height=Inches(height))
    elif width:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top))

# ============================================================================
# CHART GENERATION FUNCTIONS
# ============================================================================

def create_kpi_boxes_chart(metrics, hpeg_name, trust_avg=None):
    """Create KPI boxes for executive dashboard."""
    fig, ax = plt.subplots(1, 4, figsize=(12, 2))
    fig.patch.set_facecolor('white')

    # Data
    kpis = [
        {
            'title': 'Received (Current)',
            'value': f"{metrics['total_current']:,}",
            'change': f"{'↑' if metrics['change_absolute'] > 0 else '↓'}{abs(metrics['change_percent']):.0f}% vs prev period",
            'change_color': NHS_COLORS_HEX['nhs_dark_red'] if metrics['change_absolute'] > 0 else NHS_COLORS_HEX['nhs_green']
        },
        {
            'title': 'Ongoing (All)',
            'value': f"{metrics['ongoing']:,}",
            'change': "total active cases",
            'change_color': NHS_COLORS_HEX['nhs_mid_grey']
        },
        {
            'title': 'Open >6 Months',
            'value': f"{metrics['six_month_cases']:,}",
            'change': f"{metrics['six_month_cases']/metrics['ongoing']*100:.0f}% of ongoing" if metrics['ongoing'] > 0 else "0%",
            'change_color': NHS_COLORS_HEX['nhs_dark_red'] if metrics['six_month_cases'] > 0 else NHS_COLORS_HEX['nhs_green']
        },
        {
            'title': 'Open >12 Months',
            'value': f"{metrics['twelve_month_cases']:,}",
            'change': f"{metrics['twelve_month_cases']/metrics['ongoing']*100:.0f}% of ongoing" if metrics['ongoing'] > 0 else "0%",
            'change_color': NHS_COLORS_HEX['nhs_dark_red'] if metrics['twelve_month_cases'] > 0 else NHS_COLORS_HEX['nhs_green']
        }
    ]

    for i, (axis, kpi) in enumerate(zip(ax, kpis)):
        axis.set_xlim(0, 1)
        axis.set_ylim(0, 1)
        axis.axis('off')

        # Background box
        rect = Rectangle((0.05, 0.1), 0.9, 0.8,
                        facecolor=NHS_COLORS_HEX['nhs_pale_grey'],
                        edgecolor=NHS_COLORS_HEX['nhs_blue'], linewidth=2)
        axis.add_patch(rect)

        # Title
        axis.text(0.5, 0.75, kpi['title'], ha='center', va='center',
                 fontsize=10, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])

        # Value
        axis.text(0.5, 0.45, kpi['value'], ha='center', va='center',
                 fontsize=18, fontweight='bold', color=NHS_COLORS_HEX['nhs_blue'])

        # Change/subtitle
        axis.text(0.5, 0.2, kpi['change'], ha='center', va='center',
                 fontsize=9, color=kpi['change_color'])

    plt.tight_layout()
    return fig

def create_6month_trend_chart(df_current, df_previous, hpeg_name):
    """Create 6-month trend chart with stacked bars showing Closed vs Ongoing status."""
    fig, ax = plt.subplots(figsize=(12, 4))
    fig.patch.set_facecolor('white')

    # Combine data
    df_combined = pd.concat([df_previous, df_current])

    # Group by Month and Closed? status
    monthly_status = df_combined.groupby(['Month', 'Closed?']).size().unstack(fill_value=0)

    # Ensure we have both statuses
    if 'Yes' not in monthly_status.columns:
        monthly_status['Yes'] = 0
    if 'Ongoing' not in monthly_status.columns:
        monthly_status['Ongoing'] = 0

    # Sort by month
    monthly_status = monthly_status.sort_index()

    # Create stacked bars with NHS official colors
    x_pos = np.arange(len(monthly_status))

    # Closed cases (bottom of stack)
    bars_closed = ax.bar(x_pos, monthly_status['Yes'],
                         label='Closed',
                         color='#009639',  # NHS Green
                         edgecolor=NHS_COLORS_HEX['nhs_dark_grey'], linewidth=2)

    # Ongoing cases (top of stack)
    bars_ongoing = ax.bar(x_pos, monthly_status['Ongoing'],
                          bottom=monthly_status['Yes'],
                          label='Ongoing',
                          color='#AE2573',  # NHS Pink
                          edgecolor=NHS_COLORS_HEX['nhs_dark_grey'], linewidth=2)

    # Add value labels on each segment
    for i, (closed, ongoing) in enumerate(zip(monthly_status['Yes'], monthly_status['Ongoing'])):
        total = closed + ongoing
        # Total on top (black on white background)
        ax.text(i, total + 0.5, f'{int(total)}', ha='center', va='bottom',
                fontweight='bold', fontsize=10, color=NHS_COLORS_HEX['nhs_black'])
        # Closed count (white on dark green)
        if closed > 0:
            ax.text(i, closed/2, f'{int(closed)}', ha='center', va='center',
                    fontsize=9, color='white', fontweight='bold')
        # Ongoing count (white on dark pink)
        if ongoing > 0:
            ax.text(i, closed + ongoing/2, f'{int(ongoing)}', ha='center', va='center',
                    fontsize=9, color='white', fontweight='bold')

    # Styling
    ax.set_xticks(x_pos)
    ax.set_xticklabels([m.strftime('%b %Y') for m in monthly_status.index], rotation=0, fontsize=10)
    ax.set_ylabel('Number of Complaints Received', fontsize=12, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax.set_ylim(0, monthly_status.sum(axis=1).max() * 1.15)

    # Legend at bottom center (not overlapping bars)
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.08), ncol=2,
              frameon=True, fancybox=True, fontsize=10)

    # Grid
    ax.grid(True, alpha=0.2, linestyle='-', linewidth=0.5, axis='y')
    ax.set_axisbelow(True)

    # Spines
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(NHS_COLORS_HEX['nhs_pale_grey'])
    ax.spines['bottom'].set_color(NHS_COLORS_HEX['nhs_pale_grey'])

    plt.tight_layout()
    return fig

def create_top_5_horizontal_chart(data_dict, title, color=NHS_COLORS_HEX['nhs_blue']):
    """Create horizontal bar chart for top 5 items."""
    fig, ax = plt.subplots(figsize=(6, 4.5))
    fig.patch.set_facecolor('white')

    # Get top 5
    if len(data_dict) == 0:
        # Empty chart
        ax.text(0.5, 0.5, 'No data available', ha='center', va='center', fontsize=12)
        ax.axis('off')
        plt.tight_layout()
        return fig

    sorted_items = sorted(data_dict.items(), key=lambda x: x[1], reverse=True)[:5]
    labels = [item[0] for item in reversed(sorted_items)]  # Reverse for bottom-to-top
    values = [item[1] for item in reversed(sorted_items)]

    total = sum(data_dict.values())

    # Create bars
    y_pos = np.arange(len(labels))
    bars = ax.barh(y_pos, values, color=color, edgecolor='white', linewidth=2, alpha=0.9)

    # Add value labels
    for i, (bar, val) in enumerate(zip(bars, values)):
        width = bar.get_width()
        pct = val/total*100 if total > 0 else 0
        ax.text(width/2, bar.get_y() + bar.get_height()/2,
               f'{int(val)}', ha='center', va='center',
               color='white', fontweight='bold', fontsize=10)
        # Percentage on right
        ax.text(width + ax.get_xlim()[1]*0.02, bar.get_y() + bar.get_height()/2,
               f'{pct:.1f}%', ha='left', va='center',
               color=NHS_COLORS_HEX['nhs_mid_grey'], fontsize=9)

    # Styling
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=10)
    ax.set_xlabel('Number of Cases', fontsize=11, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax.set_title(title, fontsize=12, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_blue'], pad=15)

    ax.grid(True, alpha=0.2, axis='x')
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(NHS_COLORS_HEX['nhs_pale_grey'])
    ax.spines['bottom'].set_color(NHS_COLORS_HEX['nhs_pale_grey'])

    plt.tight_layout()
    return fig

def create_stacked_monthly_chart(df_current, column_name, title):
    """Create stacked horizontal bar chart showing top 5 items by month."""
    fig, ax = plt.subplots(figsize=(10, 5))  # Larger for better visibility
    fig.patch.set_facecolor('white')

    if len(df_current) == 0 or column_name not in df_current.columns:
        ax.text(0.5, 0.5, 'No data available', ha='center', va='center', fontsize=10)
        ax.axis('off')
        plt.tight_layout()
        return fig

    # Get overall top 5 (based on total volume)
    top_5_raw = df_current[column_name].value_counts().head(5).index.tolist()

    # Truncate long labels to 35 characters
    top_5 = []
    for label in top_5_raw:
        if isinstance(label, str) and len(label) > 35:
            top_5.append(label[:32] + '...')
        else:
            top_5.append(label)

    # Get months (should be 3 months)
    months = sorted(df_current['Month'].unique())
    month_labels = [m.strftime('%b') for m in months]

    # Build data matrix: rows = top 5 items, columns = months
    data_matrix = []
    for raw_label, display_label in zip(reversed(top_5_raw), reversed(top_5)):
        month_counts = []
        for month in months:
            count = len(df_current[(df_current[column_name] == raw_label) & (df_current['Month'] == month)])
            month_counts.append(count)
        data_matrix.append(month_counts)

    # Create stacked horizontal bars
    y_pos = np.arange(len(top_5))
    colors = [NHS_COLORS_HEX['nhs_bright_blue'], NHS_COLORS_HEX['nhs_blue'], NHS_COLORS_HEX['nhs_dark_blue']]

    left_offset = np.zeros(len(top_5))
    for month_idx, month_label in enumerate(month_labels):
        values = [row[month_idx] for row in data_matrix]
        bars = ax.barh(y_pos, values, left=left_offset, color=colors[month_idx % len(colors)],
                      label=month_label, edgecolor='white', linewidth=1)

        # Add value labels for non-zero values
        for i, (bar, val) in enumerate(zip(bars, values)):
            if val > 0:
                ax.text(left_offset[i] + val/2, bar.get_y() + bar.get_height()/2,
                       str(int(val)), ha='center', va='center',
                       color='white', fontweight='bold', fontsize=7)

        left_offset += np.array(values)

    # Add total labels
    for i, total in enumerate(left_offset):
        ax.text(total + ax.get_xlim()[1]*0.01, i,
               f'{int(total)}', ha='left', va='center',
               fontweight='bold', fontsize=9, color=NHS_COLORS_HEX['nhs_black'])

    # Styling
    ax.set_yticks(y_pos)
    ax.set_yticklabels([item for item in reversed(top_5)], fontsize=8)
    ax.set_xlabel('Number of Cases', fontsize=9, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax.set_title(title, fontsize=10, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_blue'], pad=8)

    # Legend
    ax.legend(loc='lower right', frameon=True, fancybox=True, fontsize=8, title='Month')

    ax.grid(True, alpha=0.2, axis='x')
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(NHS_COLORS_HEX['nhs_pale_grey'])
    ax.spines['bottom'].set_color(NHS_COLORS_HEX['nhs_pale_grey'])

    plt.tight_layout()
    return fig

def create_rolling_subjects_chart(subject_monthly, df_current):
    """Create grouped bar chart showing subject changes across 3 months."""
    fig, ax = plt.subplots(figsize=(12, 5.5))  # Larger for better visibility
    fig.patch.set_facecolor('white')

    if len(subject_monthly) == 0:
        ax.text(0.5, 0.5, 'No subject data available', ha='center', va='center', fontsize=10)
        ax.axis('off')
        plt.tight_layout()
        return fig

    # Get all months in order
    all_months = sorted(df_current['Month'].unique())
    month_labels = [m.strftime('%b %Y') for m in all_months]

    # Sort subjects by total volume
    subject_totals = {subj: sum(counts.values()) for subj, counts in subject_monthly.items()}
    sorted_subjects = sorted(subject_totals.items(), key=lambda x: x[1], reverse=True)

    # Prepare data matrix with truncated labels
    data_matrix = []
    subject_labels = []
    for subj, total in sorted_subjects:
        counts_by_month = [subject_monthly[subj].get(month, 0) for month in all_months]
        data_matrix.append(counts_by_month)

        # Truncate long subject names to 40 characters
        if isinstance(subj, str) and len(subj) > 40:
            subject_labels.append(subj[:37] + '...')
        else:
            subject_labels.append(subj)

    # Create grouped bars
    x = np.arange(len(subject_labels))
    width = 0.25
    colors = [NHS_COLORS_HEX['nhs_bright_blue'], NHS_COLORS_HEX['nhs_blue'], NHS_COLORS_HEX['nhs_dark_blue']]

    for i, month_idx in enumerate(range(len(all_months))):
        values = [row[month_idx] for row in data_matrix]
        offset = (i - 1) * width
        bars = ax.bar(x + offset, values, width, label=month_labels[month_idx],
                     color=colors[i % len(colors)], edgecolor='white', linewidth=1, alpha=0.9)

        # Add value labels on bars
        for bar, val in zip(bars, values):
            if val > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
                       str(int(val)), ha='center', va='bottom', fontsize=8, fontweight='bold')

    # Styling
    ax.set_xticks(x)
    ax.set_xticklabels(subject_labels, rotation=45, ha='right', fontsize=8)
    ax.set_ylabel('Number of Cases', fontsize=10, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax.set_title('Subject Trends Across 3 Months', fontsize=11, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_dark_blue'], pad=10)

    # Move legend to upper right corner within plot
    ax.legend(loc='upper right', frameon=True, fancybox=True, fontsize=8)

    ax.grid(True, alpha=0.2, axis='y')
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # Adjust layout to prevent label cutoff
    plt.tight_layout(rect=[0, 0.03, 1, 1])
    return fig

def create_topic_analysis_chart(topics, hpeg_dist, trust_avg_dist):
    """Create topic analysis visualization with deviation indicators."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    fig.patch.set_facecolor('white')

    # Left: Topic keywords
    ax1.axis('off')
    y_start = 0.95
    y_step = 0.18

    ax1.text(0.5, 0.98, 'Top 5 Topics', ha='center', va='top',
            fontsize=14, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_blue'])

    # Get top 5 topics by weight
    top_5_indices = np.argsort(hpeg_dist)[-5:][::-1]

    for i, topic_idx in enumerate(top_5_indices):
        topic = topics[topic_idx]
        y_pos = y_start - (i + 1) * y_step

        # Topic label
        ax1.text(0.05, y_pos, f"Topic {topic['id']}: {topic['label']}",
                ha='left', va='top', fontsize=11, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_blue'])

        # Keywords
        keywords_text = ', '.join(topic['keywords'])
        ax1.text(0.05, y_pos - 0.04, keywords_text,
                ha='left', va='top', fontsize=9, style='italic',
                color=NHS_COLORS_HEX['nhs_mid_grey'])

    # Right: Deviation indicators
    ax2.axis('off')
    ax2.text(0.5, 0.98, 'Change Indicators', ha='center', va='top',
            fontsize=14, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_blue'])

    # Calculate deviations
    for i, topic_idx in enumerate(top_5_indices):
        y_pos = y_start - (i + 1) * y_step

        hpeg_weight = hpeg_dist[topic_idx] * 100
        trust_weight = trust_avg_dist[topic_idx] * 100
        diff = hpeg_weight - trust_weight

        # Determine status
        if abs(diff) < 2:  # Within 2% is stable
            status = "—"
            status_color = NHS_COLORS_HEX['nhs_mid_grey']
            status_text = "Stable"
        elif diff > 0:
            status = "↑"
            status_color = NHS_COLORS_HEX['nhs_dark_red']
            status_text = f"Above average (+{diff:.1f}%)"
        else:
            status = "↓"
            status_color = NHS_COLORS_HEX['nhs_green']
            status_text = f"Below average ({diff:.1f}%)"

        # Display
        ax2.text(0.1, y_pos, status, ha='center', va='top',
                fontsize=20, fontweight='bold', color=status_color)
        ax2.text(0.25, y_pos, status_text, ha='left', va='top',
                fontsize=10, color=status_color, fontweight='bold')
        ax2.text(0.25, y_pos - 0.04, f"HPEG: {hpeg_weight:.1f}% | Trust: {trust_weight:.1f}%",
                ha='left', va='top', fontsize=9, color=NHS_COLORS_HEX['nhs_mid_grey'])

    plt.tight_layout()
    return fig

def create_topic_intelligence_chart(topic_priorities):
    """
    Create ENHANCED strategic Topic Intelligence chart with actionable insights.

    PHASES 1-4 Integration:
    - Phase 1: Trend indicators (complaint counts, up/down/stable arrows)
    - Phase 2: Team ownership (primary + secondary leads)
    - Phase 3: Anonymized examples (real complaint patterns)
    - Phase 4: Action plans (specific timelines and steps)

    4-Column Layout:
    1. Priority (Enhanced): Level, score, trend, complaint count
    2. Theme & Context: Label, keywords, prevalence, resolution
    3. Root Cause: Team ownership, anonymized examples
    4. Action Plan: Timeline, specific steps, success metric

    Args:
        topic_priorities: List of enhanced topic priority dicts

    Returns:
        matplotlib Figure with actionable topic intelligence
    """
    import matplotlib.pyplot as plt
    import numpy as np

    # NHS Colors (assuming these are defined globally)
    NHS_COLORS_HEX = {
        'nhs_blue': '#005EB8',
        'nhs_dark_blue': '#003087',
        'nhs_dark_grey': '#425563',
        'nhs_mid_grey': '#768692',
        'nhs_pale_grey': '#E8EDEE',
        'nhs_green': '#009639',
        'nhs_dark_red': '#8A1538',
        'nhs_pink': '#AE2573'
    }

    fig = plt.figure(figsize=(14, 7))  # Wider for 4 columns
    fig.patch.set_facecolor('white')

    # Create grid: 4 columns with adjusted spacing
    gs = fig.add_gridspec(6, 4, hspace=0.4, wspace=0.12,
                          left=0.05, right=0.98, top=0.90, bottom=0.08)

    # Title
    fig.text(0.5, 0.96, 'Topic Intelligence & Priority Action Plans',
             ha='center', fontsize=16, fontweight='bold',
             color=NHS_COLORS_HEX['nhs_dark_blue'])

    # Column headers
    fig.text(0.10, 0.90, 'Priority & Trend', ha='center', fontsize=10, fontweight='bold',
             color=NHS_COLORS_HEX['nhs_dark_grey'])
    fig.text(0.28, 0.90, 'Theme & Context', ha='center', fontsize=10, fontweight='bold',
             color=NHS_COLORS_HEX['nhs_dark_grey'])
    fig.text(0.50, 0.90, 'Root Cause', ha='center', fontsize=10, fontweight='bold',
             color=NHS_COLORS_HEX['nhs_dark_grey'])
    fig.text(0.75, 0.90, 'Action Plan', ha='center', fontsize=10, fontweight='bold',
             color=NHS_COLORS_HEX['nhs_dark_grey'])

    # Take top 5 priorities
    sorted_topics = sorted(topic_priorities, key=lambda x: x['priority_score'], reverse=True)[:5]

    # Y-positions for 5 topics
    y_positions = [0.75, 0.59, 0.43, 0.27, 0.11]
    row_height = 0.14

    for idx, (topic, y_pos) in enumerate(zip(sorted_topics, y_positions)):
        # ==== COLUMN 1: PRIORITY & TREND ====
        priority_level = topic['priority_level']
        priority_color = topic['priority_color']
        priority_score = topic['priority_score']

        # Phase 1 data
        trend_direction = topic.get('trend_direction', '→')
        absolute_count = topic.get('absolute_count', 0)
        trend_percentage = topic.get('trend_percentage', 0)

        # Priority box
        fig.patches.extend([plt.Rectangle((0.05, y_pos), 0.10, 0.09,
                                         facecolor=priority_color,
                                         edgecolor=NHS_COLORS_HEX['nhs_dark_grey'],
                                         linewidth=2, transform=fig.transFigure, zorder=2)])

        # Priority level + trend arrow
        fig.text(0.10, y_pos + 0.07, f"{priority_level} {trend_direction}",
                ha='center', va='center', fontsize=10, fontweight='bold', color='white')

        # Score
        fig.text(0.10, y_pos + 0.04, f"Score: {priority_score:.2f}",
                ha='center', va='center', fontsize=7, color='white')

        # Complaint count + trend %
        trend_sign = '+' if trend_percentage >= 0 else ''
        fig.text(0.10, y_pos + 0.01, f"{absolute_count} complaints",
                ha='center', va='center', fontsize=7, color='white')
        if trend_percentage != 0:
            fig.text(0.10, y_pos - 0.01, f"({trend_sign}{trend_percentage:.0f}%)",
                    ha='center', va='center', fontsize=6, color='white', style='italic')

        # ==== COLUMN 2: THEME & CONTEXT ====
        topic_label = topic['topic_label']
        keywords = ', '.join(topic['keywords'][:4])  # Top 4 keywords
        prevalence = topic.get('prevalence_pct', 0)
        median_days = topic.get('median_resolution_days', 0)

        # Topic label
        fig.text(0.17, y_pos + 0.08, topic_label,
                ha='left', va='top', fontsize=9, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_blue'], wrap=True)

        # Keywords
        fig.text(0.17, y_pos + 0.05, f"Keywords: {keywords}",
                ha='left', va='top', fontsize=7, style='italic',
                color=NHS_COLORS_HEX['nhs_mid_grey'])

        # Metrics
        fig.text(0.17, y_pos + 0.02, f"Prevalence: {prevalence:.1f}%",
                ha='left', va='top', fontsize=7, color=NHS_COLORS_HEX['nhs_dark_grey'])
        fig.text(0.17, y_pos - 0.01, f"Median Res: {median_days:.0f} days",
                ha='left', va='top', fontsize=7, color=NHS_COLORS_HEX['nhs_dark_grey'])

        # ==== COLUMN 3: ROOT CAUSE ====
        # Phase 2: Team ownership
        primary_lead = topic.get('primary_lead', 'HPEG Lead')
        primary_pct = topic.get('primary_lead_pct', 0)
        secondary_leads = topic.get('secondary_leads', [])

        # Lead contact
        fig.text(0.38, y_pos + 0.08, f"Lead: {primary_lead}",
                ha='left', va='top', fontsize=8, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_dark_blue'])
        if primary_pct > 0:
            fig.text(0.38, y_pos + 0.05, f"({primary_pct:.0f}% of complaints)",
                    ha='left', va='top', fontsize=6, color=NHS_COLORS_HEX['nhs_mid_grey'])

        # Secondary teams
        if len(secondary_leads) > 0:
            other_teams = ', '.join([f"{team[0]}" for team in secondary_leads[:2]])
            fig.text(0.38, y_pos + 0.03, f"Also: {other_teams}",
                    ha='left', va='top', fontsize=6, color=NHS_COLORS_HEX['nhs_mid_grey'])

        # Phase 3: Anonymized examples (show 1-2)
        example_complaints = topic.get('example_complaints', [])
        if len(example_complaints) > 0:
            fig.text(0.38, y_pos, "Examples:",
                    ha='left', va='top', fontsize=7, fontweight='bold',
                    color=NHS_COLORS_HEX['nhs_dark_grey'])

            for ex_idx, example in enumerate(example_complaints[:2]):  # Max 2 examples
                # Truncate example if too long
                example_short = example[:50] + '...' if len(example) > 50 else example
                fig.text(0.38, y_pos - 0.02 - (ex_idx * 0.02), f"- {example_short}",
                        ha='left', va='top', fontsize=6,
                        color=NHS_COLORS_HEX['nhs_dark_grey'])

        # ==== COLUMN 4: ACTION PLAN ====
        # Phase 4: Action timeline and steps
        action_timeline = topic.get('action_timeline', 'Review')
        action_steps = topic.get('action_steps', [])
        success_metric = topic.get('success_metric', '')

        # Timeline header
        timeline_text = f"Timeline: {action_timeline}"
        fig.text(0.60, y_pos + 0.08, timeline_text,
                ha='left', va='top', fontsize=8, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_dark_blue'])

        # Action steps (show first 3)
        step_y = y_pos + 0.05
        for step_idx, step in enumerate(action_steps[:3]):
            # Truncate step if too long
            step_short = step[:45] + '...' if len(step) > 45 else step
            fig.text(0.60, step_y - (step_idx * 0.02), f"- {step_short}",
                    ha='left', va='top', fontsize=6,
                    color=NHS_COLORS_HEX['nhs_dark_grey'])

        # Success metric
        if success_metric:
            metric_short = success_metric[:40] + '...' if len(success_metric) > 40 else success_metric
            fig.text(0.60, y_pos - 0.01, f"Target: {metric_short}",
                    ha='left', va='top', fontsize=6, style='italic',
                    color=NHS_COLORS_HEX['nhs_pink'])

        # Separator line between topics (except last)
        if idx < len(sorted_topics) - 1:
            fig.patches.extend([plt.Rectangle((0.05, y_pos - 0.03), 0.93, 0.001,
                                             facecolor=NHS_COLORS_HEX['nhs_pale_grey'],
                                             transform=fig.transFigure, zorder=1)])

    # Footer notes
    fig.text(0.5, 0.03,
             'Priority Score = (Deviation × 0.5) + (Resolution Time × 0.5)  |  Trend: ↑ Growing  ↓ Declining  → Stable  |  Examples anonymized per NHS data governance',
             ha='center', fontsize=7, style='italic', color=NHS_COLORS_HEX['nhs_mid_grey'])

    return fig

def create_deadline_compliance_chart(df_current):
    """Create deadline compliance chart."""
    fig, ax = plt.subplots(figsize=(6, 4.5))
    fig.patch.set_facecolor('white')

    # Get monthly deadline status
    months = sorted(df_current['Month'].unique())

    met_pct = []
    missed_pct = []
    in_progress_pct = []

    for month in months:
        month_data = df_current[df_current['Month'] == month]
        total = len(month_data)

        if total > 0:
            met = (month_data['Deadline Status'] == 'Deadline Met').sum() / total * 100
            missed = (month_data['Deadline Status'] == 'Deadline Missed').sum() / total * 100
            in_prog = month_data['Deadline Status'].str.contains('Progress', na=False).sum() / total * 100
        else:
            met = missed = in_prog = 0

        met_pct.append(met)
        missed_pct.append(missed)
        in_progress_pct.append(in_prog)

    # Plot lines
    x_pos = np.arange(len(months))
    ax.plot(x_pos, met_pct, marker='o', linewidth=2.5, markersize=8,
           color=NHS_COLORS_HEX['nhs_green'], label='Met (%)',
           markeredgecolor='white', markeredgewidth=2)
    ax.plot(x_pos, missed_pct, marker='s', linewidth=2.5, markersize=8,
           color=NHS_COLORS_HEX['nhs_dark_red'], label='Missed (%)',
           markeredgecolor='white', markeredgewidth=2)

    # Styling
    ax.set_xticks(x_pos)
    ax.set_xticklabels([m.strftime('%b %Y') for m in months], rotation=0)
    ax.set_ylabel('Percentage', fontsize=11, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax.set_title('Deadline Compliance Trend', fontsize=12, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_dark_blue'], pad=15)
    ax.set_ylim(0, 100)

    ax.legend(loc='best', frameon=True, fancybox=True)
    ax.grid(True, alpha=0.2)
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    plt.tight_layout()
    return fig

def create_complexity_donut_chart(complexity_dist):
    """Create complexity distribution donut chart."""
    fig, ax = plt.subplots(figsize=(6, 4))
    fig.patch.set_facecolor('white')

    # Define expected categories with NHS official colors
    categories = ['Basic', 'Regular', 'Complex']
    colors = ['#009639', '#005EB8', '#AE2573']  # Basic=Green, Regular=Blue, Complex=Pink

    values = [complexity_dist.get(cat, 0) for cat in categories]
    total = sum(values)

    if total == 0:
        ax.text(0.5, 0.5, 'No complexity data', ha='center', va='center')
        ax.axis('off')
        plt.tight_layout()
        return fig

    # Create donut with darker borders for definition
    wedges, texts, autotexts = ax.pie(values, labels=categories, colors=colors,
                                       autopct='%1.0f%%', startangle=90,
                                       pctdistance=0.85,
                                       wedgeprops=dict(width=0.4, edgecolor=NHS_COLORS_HEX['nhs_dark_grey'], linewidth=3))

    for autotext in autotexts:
        autotext.set_color('white')  # White text on dark NHS colors for readability
        autotext.set_fontweight('bold')
        autotext.set_fontsize(12)

    # Center text
    ax.text(0, 0, f'{total}\nCases', ha='center', va='center',
           fontsize=14, fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_grey'])

    ax.set_title('Complexity Distribution', fontsize=12, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_dark_blue'], pad=15)

    plt.tight_layout()
    return fig

def create_response_time_chart(all_hpeg_metrics, current_hpeg_name):
    """Create response time analysis chart comparing ALL HPEGs with current HPEG highlighted."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5.5))
    fig.patch.set_facecolor('white')

    # Define HPEG order (shortened names for display)
    hpeg_order = ['BHH Exec Team', 'QEH Exec Team', 'GHH Exec Team',
                  'SH Exec Team', 'W&C Exec Team', 'CSS Exec Team']
    hpeg_labels = ['BHH', 'QEHB', 'GHH', 'SH', 'W&C', 'CSS']

    # Check if we have data for current HPEG
    current_metrics = all_hpeg_metrics.get(current_hpeg_name, {})
    if current_metrics.get('resolution_mean', 0) == 0:
        ax1.text(0.5, 0.5, 'No resolution time data available', ha='center', va='center')
        ax1.axis('off')
        ax2.text(0.5, 0.5, 'No resolution time data available', ha='center', va='center')
        ax2.axis('off')
        plt.tight_layout()
        return fig

    # Left chart: Mean vs Median for all HPEGs (grouped bars)
    mean_values = []
    median_values = []
    for hpeg in hpeg_order:
        metrics = all_hpeg_metrics.get(hpeg, {})
        mean_values.append(metrics.get('resolution_mean', 0))
        median_values.append(metrics.get('resolution_median', 0))

    x_pos = np.arange(len(hpeg_labels))
    width = 0.35

    # Determine colors (highlight current HPEG)
    current_idx = hpeg_order.index(current_hpeg_name) if current_hpeg_name in hpeg_order else -1
    colors_mean = [NHS_COLORS_HEX['nhs_dark_blue'] if i == current_idx else NHS_COLORS_HEX['nhs_blue']
                   for i in range(len(hpeg_labels))]
    colors_median = [NHS_COLORS_HEX['nhs_bright_blue'] if i == current_idx else NHS_COLORS_HEX['nhs_aqua_blue']
                     for i in range(len(hpeg_labels))]

    bars1 = ax1.bar(x_pos - width/2, mean_values, width, label='Mean',
                    color=colors_mean, edgecolor='white', linewidth=1.5, alpha=0.9)
    bars2 = ax1.bar(x_pos + width/2, median_values, width, label='Median',
                    color=colors_median, edgecolor='white', linewidth=1.5, alpha=0.9)

    # Add value labels to ALL bars for clarity
    max_val = max(mean_values + median_values) if mean_values + median_values else 1
    for i, (mean_val, median_val) in enumerate(zip(mean_values, median_values)):
        if mean_val > 0:  # Only label if there's data
            ax1.text(i - width/2, mean_val + max_val*0.02,
                    f'{mean_val:.1f}', ha='center', va='bottom',
                    fontsize=8, fontweight='bold' if i == current_idx else 'normal')
        if median_val > 0:  # Only label if there's data
            ax1.text(i + width/2, median_val + max_val*0.02,
                    f'{median_val:.1f}', ha='center', va='bottom',
                    fontsize=8, fontweight='bold' if i == current_idx else 'normal')

    # Styling for left chart
    ax1.set_xticks(x_pos)
    ax1.set_xticklabels(hpeg_labels, fontsize=10)
    ax1.set_ylabel('Days to Close (Business Days)', fontsize=10, fontweight='bold',
                   color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax1.set_title('Resolution Time Comparison Across HPEGs', fontsize=11, fontweight='bold',
                  color=NHS_COLORS_HEX['nhs_dark_blue'], pad=10)
    ax1.legend(loc='upper left', fontsize=9, frameon=True)
    ax1.grid(True, alpha=0.2, axis='y')
    ax1.set_axisbelow(True)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)

    # Right chart: Target compliance by complexity for all HPEGs (grouped bars)
    complexities = ['Basic', 'Regular', 'Complex']
    targets = [25, 40, 65]

    # Get compliance data for all HPEGs
    compliance_data = {comp: [] for comp in complexities}
    for hpeg in hpeg_order:
        metrics = all_hpeg_metrics.get(hpeg, {})
        targets_met = metrics.get('targets_met_by_complexity', {})
        for comp in complexities:
            compliance_data[comp].append(targets_met.get(comp, 0))

    # Create grouped bars with NHS official colors
    x_pos2 = np.arange(len(hpeg_labels))
    width2 = 0.25
    colors_complex = ['#009639', '#005EB8', '#AE2573']  # Basic=Green, Regular=Blue, Complex=Pink

    for idx, (comp, color) in enumerate(zip(complexities, colors_complex)):
        offset = (idx - 1) * width2
        values = compliance_data[comp]
        # All bars same color, darker border for current HPEG
        bar_colors = [color for i in range(len(hpeg_labels))]
        border_colors = [NHS_COLORS_HEX['nhs_black'] if i == current_idx else NHS_COLORS_HEX['nhs_mid_grey']
                         for i in range(len(hpeg_labels))]
        border_widths = [3 if i == current_idx else 1.5 for i in range(len(hpeg_labels))]

        bars = ax2.bar(x_pos2 + offset, values, width2, label=f'{comp} ({targets[idx]}d)',
                      color=bar_colors)

        # Apply individual border styling and add value labels
        for i, (bar, val) in enumerate(zip(bars, values)):
            bar.set_edgecolor(border_colors[i])
            bar.set_linewidth(border_widths[i])
            # Add label if value > 0 (only show percentages where there's data)
            if val > 0:
                ax2.text(bar.get_x() + bar.get_width()/2, val + 2,
                        f'{val:.0f}%', ha='center', va='bottom',
                        fontsize=8, fontweight='bold' if i == current_idx else 'normal',
                        color=NHS_COLORS_HEX['nhs_black'])

    # Styling for right chart
    ax2.set_xticks(x_pos2)
    ax2.set_xticklabels(hpeg_labels, fontsize=10)
    ax2.set_ylabel('% Meeting Target', fontsize=10, fontweight='bold',
                   color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax2.set_title('Target Compliance by Complexity Level', fontsize=11, fontweight='bold',
                  color=NHS_COLORS_HEX['nhs_dark_blue'], pad=10)
    ax2.set_ylim(0, 105)
    ax2.legend(loc='upper left', fontsize=9, frameon=True, ncol=1)
    ax2.grid(True, alpha=0.2, axis='y')
    ax2.set_axisbelow(True)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)

    plt.tight_layout()
    return fig

def create_12month_trends_chart(trend_data, hpeg_items, data_type='cdg', title='12-Month Resolution Trends'):
    """Create 12-month median resolution time trends as HEAT MAP for better readability."""
    fig, ax = plt.subplots(figsize=(12, 6))
    fig.patch.set_facecolor('white')

    if data_type == 'cdg':
        trends_dict = trend_data.get('resolution_by_cdg_monthly', {})
    else:  # specialty
        trends_dict = trend_data.get('resolution_by_specialty_monthly', {})

    # FILTER: Only include items that belong to this HPEG
    trends_dict = {item: months for item, months in trends_dict.items() if item in hpeg_items}

    if len(trends_dict) == 0:
        ax.text(0.5, 0.5, 'No trend data available for this HPEG', ha='center', va='center')
        ax.axis('off')
        plt.tight_layout()
        return fig

    # Convert Period objects to datetime for plotting
    all_months = set()
    for item_trends in trends_dict.values():
        all_months.update(item_trends.keys())

    # FILTER: Only show data from December 2024 onwards, limit to 12 months
    cutoff_date = pd.Period('2024-12', 'M')
    sorted_months = sorted([m for m in list(all_months) if pd.Period(m) >= cutoff_date])
    sorted_months = sorted_months[-12:]  # Limit to last 12 months

    if len(sorted_months) == 0:
        ax.text(0.5, 0.5, 'No data available from December 2024 onwards', ha='center', va='center')
        ax.axis('off')
        plt.tight_layout()
        return fig

    month_labels = [pd.Period(m).strftime('%b %y') if isinstance(m, pd.Period) else str(m) for m in sorted_months]

    # Get top 8 items by total resolution time
    item_totals = {item: sum(values.get(m, 0) for m in sorted_months) for item, values in trends_dict.items()}
    top_items = sorted(item_totals.items(), key=lambda x: x[1], reverse=True)[:8]
    item_names = [item for item, _ in top_items]

    # Build heat map data matrix
    heat_data = []
    for item in item_names:
        row = [trends_dict[item].get(month, np.nan) for month in sorted_months]
        heat_data.append(row)

    heat_data = np.array(heat_data)

    # Create heat map with custom colormap (green=good, yellow=warning, red=bad)
    from matplotlib.colors import LinearSegmentedColormap
    colors_map = ['#009639', '#78BE20', '#FFB81C', '#FF8C00', '#8A1538']  # Green to Red
    n_bins = 100
    cmap = LinearSegmentedColormap.from_list('custom', colors_map, N=n_bins)

    # Plot heat map
    im = ax.imshow(heat_data, aspect='auto', cmap=cmap, vmin=0, vmax=70)

    # Set ticks and labels
    ax.set_xticks(np.arange(len(month_labels)))
    ax.set_yticks(np.arange(len(item_names)))
    ax.set_xticklabels(month_labels, fontsize=9)
    ax.set_yticklabels([str(name)[:30] for name in item_names], fontsize=9)

    # Rotate x labels
    plt.setp(ax.get_xticklabels(), rotation=45, ha='right', rotation_mode='anchor')

    # Add value annotations
    for i in range(len(item_names)):
        for j in range(len(month_labels)):
            value = heat_data[i, j]
            if not np.isnan(value):
                text_color = 'white' if value > 35 else 'black'
                ax.text(j, i, f'{value:.0f}', ha='center', va='center',
                       fontsize=8, fontweight='bold', color=text_color)

    # Add colorbar with target reference lines
    cbar = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label('Median Days to Close', fontsize=10, fontweight='bold')

    # Add target lines to colorbar
    cbar.ax.axhline(y=25/70, color='white', linestyle='--', linewidth=2, alpha=0.8)
    cbar.ax.axhline(y=40/70, color='white', linestyle='--', linewidth=2, alpha=0.8)
    cbar.ax.axhline(y=65/70, color='white', linestyle='--', linewidth=2, alpha=0.8)
    cbar.ax.text(1.5, 25/70, 'Basic', fontsize=7, va='center', color='black', fontweight='bold')
    cbar.ax.text(1.5, 40/70, 'Regular', fontsize=7, va='center', color='black', fontweight='bold')
    cbar.ax.text(1.5, 65/70, 'Complex', fontsize=7, va='center', color='black', fontweight='bold')

    # Title
    ax.set_title(title, fontsize=12, fontweight='bold',
                color=NHS_COLORS_HEX['nhs_dark_blue'], pad=15)

    plt.tight_layout()
    return fig

def create_seasonal_volume_chart(volume_monthly):
    """Create seasonal volume pattern chart showing RECEIVED COMPLAINTS over 12 months."""
    fig, ax = plt.subplots(figsize=(12, 5))
    fig.patch.set_facecolor('white')

    if len(volume_monthly) == 0:
        ax.text(0.5, 0.5, 'No volume data available', ha='center', va='center')
        ax.axis('off')
        plt.tight_layout()
        return fig

    # FILTER: Only show data from December 2024 onwards, limit to 12 months
    cutoff_date = pd.Period('2024-12', 'M')
    filtered_data = {m: v for m, v in volume_monthly.items() if pd.Period(m) >= cutoff_date}

    if len(filtered_data) == 0:
        ax.text(0.5, 0.5, 'No data available from December 2024 onwards', ha='center', va='center')
        ax.axis('off')
        plt.tight_layout()
        return fig

    # Sort by month and limit to last 12 months
    sorted_months = sorted(filtered_data.keys())
    sorted_months = sorted_months[-12:]  # Limit to exactly 12 months
    month_labels = [pd.Period(m).strftime('%b %Y') if isinstance(m, pd.Period) else str(m)
                    for m in sorted_months]
    volumes = [filtered_data[m] for m in sorted_months]

    # Create bar chart
    x_pos = np.arange(len(month_labels))
    bars = ax.bar(x_pos, volumes, color=NHS_COLORS_HEX['nhs_blue'],
                  edgecolor='white', linewidth=2, alpha=0.9)

    # Add value labels
    for bar, vol in zip(bars, volumes):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, height + max(volumes)*0.01,
               str(int(vol)), ha='center', va='bottom',
               fontsize=9, fontweight='bold')

    # Add trend line (only if we have at least 2 data points)
    if len(volumes) >= 2:
        z = np.polyfit(x_pos, volumes, 1)
        p = np.poly1d(z)
        ax.plot(x_pos, p(x_pos), color=NHS_COLORS_HEX['nhs_dark_red'],
               linestyle='--', linewidth=2, alpha=0.7, label='Trend Line')

    # Styling
    ax.set_xticks(x_pos)
    ax.set_xticklabels(month_labels, rotation=45, ha='right', fontsize=9)
    ax.set_ylabel('Number of Received Complaints', fontsize=10, fontweight='bold',
                  color=NHS_COLORS_HEX['nhs_dark_grey'])
    ax.set_title('Received Complaints: 12-Month Volume & Seasonal Patterns', fontsize=11,
                fontweight='bold', color=NHS_COLORS_HEX['nhs_dark_blue'], pad=10)
    ax.set_ylim(0, max(volumes) * 1.15)

    if len(volumes) >= 2:
        ax.legend(loc='upper right', fontsize=8)
    ax.grid(True, alpha=0.2, axis='y')
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    plt.tight_layout()
    return fig

def create_subject_deepdive_chart(df_current, top_n=3):
    """Create subject deep-dive analysis showing top subjects with CDG breakdown."""
    fig, axes = plt.subplots(1, top_n, figsize=(12, 5))
    fig.patch.set_facecolor('white')

    if top_n == 1:
        axes = [axes]

    # Get top N subjects
    top_subjects = df_current['Subjects'].value_counts().head(top_n).index.tolist()

    if len(top_subjects) == 0:
        axes[0].text(0.5, 0.5, 'No subject data available', ha='center', va='center')
        axes[0].axis('off')
        plt.tight_layout()
        return fig

    for idx, (ax, subject) in enumerate(zip(axes, top_subjects)):
        # Filter data for this subject
        subject_data = df_current[df_current['Subjects'] == subject]

        # Get top 5 CDGs for this subject
        cdg_counts = subject_data['CDG'].value_counts().head(5)

        if len(cdg_counts) == 0:
            ax.text(0.5, 0.5, 'No data', ha='center', va='center', fontsize=8)
            ax.axis('off')
            continue

        # Create horizontal bars (use standard NHS blue)
        y_pos = np.arange(len(cdg_counts))
        bars = ax.barh(y_pos, cdg_counts.values,
                      color=NHS_COLORS_HEX['nhs_blue'],
                      edgecolor='white', linewidth=1.5, alpha=0.9)

        # Add value labels
        for bar, val in zip(bars, cdg_counts.values):
            width = bar.get_width()
            ax.text(width + max(cdg_counts.values)*0.02, bar.get_y() + bar.get_height()/2,
                   str(int(val)), ha='left', va='center',
                   fontsize=8, fontweight='bold')

        # Styling
        ax.set_yticks(y_pos)
        ax.set_yticklabels(cdg_counts.index, fontsize=7)
        ax.set_xlabel('Cases', fontsize=8, fontweight='bold',
                     color=NHS_COLORS_HEX['nhs_dark_grey'])

        # Truncate long subject names for title
        title_text = subject if len(str(subject)) <= 35 else str(subject)[:32] + '...'
        ax.set_title(title_text, fontsize=9, fontweight='bold',
                    color=NHS_COLORS_HEX['nhs_dark_blue'], pad=5)

        # Add total count
        total = len(subject_data)
        ax.text(0.5, -0.15, f'Total: {total} cases', ha='center', va='top',
               transform=ax.transAxes, fontsize=7, style='italic',
               color=NHS_COLORS_HEX['nhs_mid_grey'])

        ax.grid(True, alpha=0.2, axis='x')
        ax.set_axisbelow(True)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

    plt.tight_layout()
    return fig

# ============================================================================
# SLIDE GENERATION FUNCTIONS
# ============================================================================

def create_slide_1_executive_dashboard(prs, hpeg_name, metrics, temp_dir):
    """Slide 1: Executive Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_bar(slide, hpeg_name, "Executive Dashboard")

    # Create KPI boxes chart
    kpi_fig = create_kpi_boxes_chart(metrics, hpeg_name)
    kpi_path = temp_dir / f"kpi_{hpeg_name}.png"
    kpi_fig.savefig(kpi_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(kpi_fig)

    # Add KPI boxes
    add_image_to_slide(slide, kpi_path, left=0.6, top=1.1, width=11.8)

    # Create 6-month trend chart
    trend_fig = create_6month_trend_chart(metrics['dataframe_current'],
                                         metrics['dataframe_previous'], hpeg_name)
    trend_path = temp_dir / f"trend_{hpeg_name}.png"
    trend_fig.savefig(trend_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(trend_fig)

    # Add trend chart
    add_image_to_slide(slide, trend_path, left=0.6, top=2.8, width=11.8)

    # Alert strip if needed
    if metrics['six_month_cases'] > 0 or metrics['safeguarding_cases'] > 0:
        alert_text = f"⚠ {metrics['six_month_cases']} cases open >6 months"
        if metrics['safeguarding_cases'] > 0:
            alert_text += f" | {metrics['safeguarding_cases']} safeguarding concerns"

        alert_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(6.7), Inches(12), Inches(0.6)
        )
        alert_box.fill.solid()
        alert_box.fill.fore_color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_dark_red'])
        alert_box.line.fill.background()

        text_frame = alert_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = alert_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def create_slide_2_locations(prs, hpeg_name, metrics, temp_dir):
    """Slide 2: Top 5 Locations with Monthly Breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Complaints by Location")

    # Stacked locations chart (monthly breakdown) - centered and larger
    loc_fig = create_stacked_monthly_chart(metrics['dataframe_current'], 'Location',
                                          'Top 5 Locations - Monthly Breakdown')
    loc_path = temp_dir / f"locations_{hpeg_name}.png"
    loc_fig.savefig(loc_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(loc_fig)

    # Center the chart on the slide (larger size)
    add_image_to_slide(slide, loc_path, left=1.0, top=1.3, width=11.3)

def create_slide_3_specialties(prs, hpeg_name, metrics, temp_dir):
    """Slide 3: Top 5 Specialties with Monthly Breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Complaints by Specialty")

    # Stacked specialties chart (monthly breakdown) - centered and larger
    spec_fig = create_stacked_monthly_chart(metrics['dataframe_current'], 'Specialty',
                                           'Top 5 Specialties - Monthly Breakdown')
    spec_path = temp_dir / f"specialties_{hpeg_name}.png"
    spec_fig.savefig(spec_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(spec_fig)

    # Center the chart on the slide (larger size)
    add_image_to_slide(slide, spec_path, left=1.0, top=1.3, width=11.3)

def create_slide_cdg_breakdown(prs, hpeg_name, metrics, temp_dir):
    """Slide 4: CDG Breakdown with Monthly Trends."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Clinical Decision Group (CDG) Breakdown")

    # Stacked CDG chart (monthly breakdown) - centered on slide
    cdg_fig = create_stacked_monthly_chart(metrics['dataframe_current'], 'CDG',
                                          'Top 5 CDGs - Monthly Breakdown')
    cdg_path = temp_dir / f"cdgs_{hpeg_name}.png"
    cdg_fig.savefig(cdg_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(cdg_fig)

    # Center the chart on the slide (larger size)
    add_image_to_slide(slide, cdg_path, left=1.0, top=1.3, width=11.3)

def create_slide_3_rolling_subjects(prs, hpeg_name, metrics, temp_dir):
    """Slide 4: Rolling Subject Analysis (formerly Slide 3)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "What's Changing? - Subject Trends")

    # Create chart
    subj_fig = create_rolling_subjects_chart(metrics['subject_monthly'],
                                             metrics['dataframe_current'])
    subj_path = temp_dir / f"subjects_{hpeg_name}.png"
    subj_fig.savefig(subj_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(subj_fig)

    # Chart spans nearly full width
    add_image_to_slide(slide, subj_path, left=0.6, top=1.2, width=12.1)

def create_slide_narrative_insights(prs, hpeg_name, metrics):
    """Slide: Narrative Insights - Key Trends & Patterns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Intelligence Report: Key Trends & Patterns")

    insights = metrics.get('narrative_insights', [])

    if len(insights) == 0:
        # No significant insights
        no_data_box = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(9.33), Inches(2)
        )
        text_frame = no_data_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = "No significant trends detected for this period."
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_mid_grey'])
        p.alignment = PP_ALIGN.CENTER
        return

    # Executive Summary Box at top
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.85)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_bright_blue'])
    summary_box.line.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_blue'])
    summary_box.line.width = Pt(2)

    text_frame = summary_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)

    # Count increases vs decreases
    increases = [i for i in insights if i['change'] > 0]
    decreases = [i for i in insights if i['change'] < 0]
    high_priority = [i for i in insights if i['priority'] == 'high']

    # Summary text
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Title
    run = p.add_run()
    run.text = "EXECUTIVE SUMMARY: "
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Get period info from first insight if available
    period_text = ""
    if len(insights) > 0 and 'text' in insights[0]:
        # Extract month from narrative (e.g., "November 2025")
        import re
        match = re.search(r'in ([A-Z][a-z]+ \d{4})', insights[0]['text'])
        if match:
            latest_month = match.group(1)
            period_text = f"Latest month ({latest_month}) vs previous 2-month average"

    # Summary stats
    run = p.add_run()
    summary_text = f"{len(insights)} significant trends identified"
    if len(high_priority) > 0:
        summary_text += f" ({len(high_priority)} high priority)"
    if len(increases) > 0 and len(decreases) > 0:
        summary_text += f" • {len(increases)} increases, {len(decreases)} decreases"
    elif len(increases) > 0:
        summary_text += f" • {len(increases)} increases"
    elif len(decreases) > 0:
        summary_text += f" • {len(decreases)} decreases"

    if period_text:
        summary_text += f" • {period_text}"

    run.text = summary_text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Separate increases and decreases
    top_section = 2.3
    left_margin = 0.8
    column_width = 11.7

    # Section headers and insights
    current_y = top_section

    # INCREASES section
    if len(increases) > 0:
        # Section header
        header_box = slide.shapes.add_textbox(
            Inches(left_margin), Inches(current_y),
            Inches(column_width), Inches(0.35)
        )
        header_frame = header_box.text_frame
        p = header_frame.paragraphs[0]
        run = p.add_run()
        run.text = "▲ AREAS OF INCREASE"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_dark_red'])

        current_y += 0.45

        # Add increase insights
        for insight in increases[:4]:  # Top 4 increases
            insight_box = _create_insight_card(slide, insight, left_margin, current_y, column_width)
            current_y += 0.85

    # DECREASES section
    if len(decreases) > 0:
        # Add spacing if we had increases
        if len(increases) > 0:
            current_y += 0.15

        # Section header
        header_box = slide.shapes.add_textbox(
            Inches(left_margin), Inches(current_y),
            Inches(column_width), Inches(0.35)
        )
        header_frame = header_box.text_frame
        p = header_frame.paragraphs[0]
        run = p.add_run()
        run.text = "▼ AREAS OF DECREASE"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_green'])

        current_y += 0.45

        # Add decrease insights
        for insight in decreases[:2]:  # Top 2 decreases
            insight_box = _create_insight_card(slide, insight, left_margin, current_y, column_width)
            current_y += 0.95

def _create_insight_card(slide, insight, left, top, width):
    """Create a professional insight card with NHS branding."""
    # Determine colors based on priority and direction
    if insight['change'] > 0:
        # Increase
        if insight['priority'] == 'high':
            border_color = NHS_COLORS_RGB['nhs_dark_red']
            bg_color = (255, 245, 245)  # Very light red
            priority_text = "HIGH PRIORITY"
        elif insight['priority'] == 'medium':
            border_color = NHS_COLORS_RGB['nhs_orange']
            bg_color = (255, 250, 240)  # Very light orange
            priority_text = "MEDIUM"
        else:
            border_color = NHS_COLORS_RGB['nhs_warm_yellow']
            bg_color = (255, 255, 245)  # Very light yellow
            priority_text = "MONITOR"
    else:
        # Decrease
        if insight['priority'] == 'high':
            border_color = NHS_COLORS_RGB['nhs_green']
            bg_color = (245, 255, 245)  # Very light green
            priority_text = "SIGNIFICANT"
        else:
            border_color = NHS_COLORS_RGB['nhs_mid_grey']
            bg_color = (250, 250, 250)  # Very light grey
            priority_text = "POSITIVE"

    # Create card with border (taller to fit 3 lines)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.85)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*bg_color)
    card.line.color.rgb = RGBColor(*border_color)
    card.line.width = Pt(2.5)

    # Add text
    text_frame = card.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.15)
    text_frame.margin_right = Inches(0.15)
    text_frame.margin_top = Inches(0.08)
    text_frame.margin_bottom = Inches(0.08)

    # First paragraph: CDG and priority badge
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # CDG name in bold
    run = p.add_run()
    run.text = insight['cdg']
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_dark_blue'])

    # Priority badge
    run = p.add_run()
    run.text = f" [{priority_text}]"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*border_color)

    # Trend arrow and change
    run = p.add_run()
    arrow = "↑" if insight['change'] > 0 else "↓"
    change_text = f" {arrow} {abs(int(insight['change']))} cases ({abs(insight['pct_change']):.0f}%)"
    run.text = change_text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*border_color)

    # Second paragraph: The narrative (subjects)
    p2 = text_frame.add_paragraph()
    p2.level = 0
    p2.space_before = Pt(2)

    # Extract subjects from the original text
    original_text = insight['text']
    if 'relating to ' in original_text:
        subjects_part = original_text.split('relating to ')[1].rstrip('.')

        run = p2.add_run()
        run.text = f"Main themes: {subjects_part}"
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_dark_grey'])

    # Third paragraph: Action recommendation
    p3 = text_frame.add_paragraph()
    p3.level = 0
    p3.space_before = Pt(3)

    run = p3.add_run()
    run.text = "→ "
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_blue'])

    run = p3.add_run()
    run.text = insight.get('action', '')
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_dark_blue'])

    return card

def create_slide_4_topic_intelligence(prs, hpeg_name, topic_priorities, temp_dir):
    """
    Slide 4: Strategic Topic Intelligence & Priority Themes.

    Displays top 5 priority topics with:
    - Priority level (CRITICAL/MONITOR/MAINTAIN)
    - Topic keywords and metrics
    - Actionable recommendations
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Topic Intelligence & Priority Themes")

    # Check if topic priorities exist and have data
    if not topic_priorities or len(topic_priorities) == 0:
        # No topic data available - add message
        add_text_box(slide, "Insufficient data for topic analysis (minimum 10 complaints required)",
                    left=2, top=3, width=9, height=1, font_size=14,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.CENTER)
        return

    # Create strategic topic intelligence chart
    topic_fig = create_topic_intelligence_chart(topic_priorities)
    topic_path = temp_dir / f"topic_intelligence_{hpeg_name}.png"
    topic_fig.savefig(topic_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(topic_fig)

    # Add chart to slide (full width, positioned below title)
    add_image_to_slide(slide, topic_path, left=0.5, top=1.1, width=12)

def create_slide_5_performance_metrics(prs, hpeg_name, metrics, temp_dir):
    """Slide 5: Complexity Distribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Complaint Complexity Distribution")

    # Complexity chart (centered, larger)
    complex_fig = create_complexity_donut_chart(metrics['complexity_dist'])
    complex_path = temp_dir / f"complexity_{hpeg_name}.png"
    complex_fig.savefig(complex_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(complex_fig)

    # Center the chart on the slide (adjusted to prevent overlap at bottom)
    add_image_to_slide(slide, complex_path, left=3.5, top=1.5, width=6)

def create_slide_6_risk_dashboard(prs, hpeg_name, metrics):
    """Slide 6: Risk Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Risk Dashboard")

    # Long-standing cases text
    six_month_count = metrics['six_month_cases']
    add_text_box(slide, f"Cases Open >6 Months: {six_month_count}",
                left=0.5, top=1.5, width=12, height=0.5,
                font_size=16, bold=True, color_rgb=NHS_COLORS_RGB['nhs_dark_red'])

    # Safeguarding text
    sg_count = metrics['safeguarding_cases']
    if sg_count > 0:
        add_text_box(slide, f"Safeguarding Concerns: {sg_count}",
                    left=0.5, top=2.2, width=12, height=0.5,
                    font_size=16, bold=True, color_rgb=NHS_COLORS_RGB['nhs_dark_red'])

    # Note
    add_text_box(slide,
                "Detailed breakdowns available in operational dashboards.\nFocus on specialty-level resolution for long-standing cases.",
                left=0.5, top=4, width=12, height=1,
                font_size=11, color_rgb=NHS_COLORS_RGB['nhs_mid_grey'])

def create_slide_7_demographics(prs, hpeg_name, metrics, demographic_findings):
    """Slide 7: Demographic Insights (conditional)."""
    # Only create if there are findings
    if len(demographic_findings) == 0:
        return  # Skip this slide

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Demographic Insights")

    # Display findings
    y_pos = 1.5
    add_text_box(slide, "Statistically Significant Findings (p < 0.05):",
                left=0.5, top=y_pos, width=12, height=0.5,
                font_size=14, bold=True, color_rgb=NHS_COLORS_RGB['nhs_dark_blue'])

    y_pos += 0.7
    for finding in demographic_findings[:5]:  # Show top 5
        text = f"• {finding['interpretation']}"
        add_text_box(slide, text, left=1, top=y_pos, width=11, height=0.4,
                    font_size=11, color_rgb=NHS_COLORS_RGB['nhs_black'])
        y_pos += 0.5

def create_slide_8_oldest_cases(prs, hpeg_name):
    """Slide 8: 10 Oldest Cases (blank template)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "10 Oldest Cases")

    # Instruction text
    instruction = ("This slide must be completed manually by the HPEG lead.\n\n"
                  "Contains patient identifiable information - add details of the 10 oldest "
                  "ongoing cases including: ID, Date Received, Days Open, Specialty, Subject, "
                  "and brief description.\n\n"
                  "Use table format for clarity.")

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = instruction
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_mid_grey'])
    p.font.name = 'Arial'
    p.font.italic = True

def create_slide_9_actions_status(prs, hpeg_name):
    """Slide 9: Actions Status (blank template)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Actions Status")

    # Three sections with placeholder text
    sections = [
        ("Outstanding Actions from Last Report", 1.5),
        ("New Actions Identified This Quarter", 3.2),
        ("Completed Actions", 4.9)
    ]

    for section_title, y_pos in sections:
        # Section header
        header = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.4))
        header_frame = header.text_frame
        p = header_frame.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_blue'])

        # Placeholder bullets
        content = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.5), Inches(11), Inches(1.2))
        content_frame = content.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = "• [Add action item here]\n• [Add action item here]\n• [Add action item here]"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_pale_grey'])
        p.font.italic = True

def create_slide_10_current_performance(prs, hpeg_name):
    """Slide 10: Current Performance (blank template)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Current Performance & Key Achievements")

    # Two sections
    sections = [
        ("Performance Highlights", 1.5),
        ("Areas for Improvement / Focus", 4.0)
    ]

    for section_title, y_pos in sections:
        # Section header
        header = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.4))
        header_frame = header.text_frame
        p = header_frame.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_blue'])

        # Placeholder
        content = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.5), Inches(12), Inches(2))
        content_frame = content.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = "[Add narrative here - key achievements, notable improvements, strategic priorities]"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_pale_grey'])
        p.font.italic = True

def create_slide_response_time_analysis(prs, hpeg_name, all_hpeg_metrics, temp_dir):
    """NEW Slide: Response Time Analysis comparing ALL HPEGs with current HPEG highlighted."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Complaint Resolution Efficiency")

    # Create chart showing ALL HPEGs with current highlighted
    response_fig = create_response_time_chart(all_hpeg_metrics, hpeg_name)
    response_path = temp_dir / f"response_time_{hpeg_name}.png"
    response_fig.savefig(response_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
    plt.close(response_fig)

    # Add chart to slide
    add_image_to_slide(slide, response_path, left=0.6, top=1.2, width=11.8)

def create_slide_12month_cdg_trends(prs, hpeg_name, trends_12month, hpeg_cdgs, metrics, temp_dir):
    """NEW Slide: Current Period vs 12-Month Rolling Average - CDGs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Resolution Time Trends by CDG")

    # Get current period months from dataframe
    df_current = metrics.get('dataframe_current')
    if df_current is None or len(df_current) == 0:
        add_text_box(slide, "No data available",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    # Get the 3 months in the current reporting period
    if 'Completed Date' in df_current.columns:
        df_temp = df_current.copy()
        df_temp['Completed Month'] = pd.to_datetime(df_temp['Completed Date'], errors='coerce').dt.to_period('M')
        current_months = sorted(df_temp['Completed Month'].dropna().unique())[-3:]  # Last 3 months
    else:
        add_text_box(slide, "No completion date data available",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    if len(current_months) == 0:
        add_text_box(slide, "No completion data for current period",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    # Format month labels
    month_labels = [m.strftime('%b %y') for m in current_months]
    period_label = f"{month_labels[0]}-{month_labels[-1]}"

    # Analyze trend data
    cdg_trends = trends_12month.get('resolution_by_cdg_monthly', {})
    cdg_trends_filtered = {cdg: months for cdg, months in cdg_trends.items() if cdg in hpeg_cdgs}

    if len(cdg_trends_filtered) == 0:
        add_text_box(slide, "No CDG trend data available for this HPEG",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    insights = []
    for cdg, monthly_data in cdg_trends_filtered.items():
        # Get last 12 months for rolling average
        cutoff_date = pd.Period('2024-12', 'M')
        filtered_months = {m: v for m, v in monthly_data.items() if pd.Period(m) >= cutoff_date}
        sorted_months = sorted(filtered_months.keys())[-12:]  # Last 12 months
        twelve_month_values = [filtered_months[m] for m in sorted_months]

        # Calculate 12-month rolling median
        rolling_median = np.median(twelve_month_values) if len(twelve_month_values) > 0 else 0

        # Get current 3-month values
        current_period_values = [monthly_data.get(m, None) for m in current_months]
        current_period_values = [v for v in current_period_values if v is not None and v > 0]

        if len(current_period_values) == 0:
            continue

        # Calculate current period median
        current_median = np.median(current_period_values)

        # Calculate difference
        trend = current_median - rolling_median

        # Determine performance
        if current_median <= 25:
            performance = "Excellent (within Basic target)"
        elif current_median <= 40:
            performance = "Good (within Regular target)"
        elif current_median <= 65:
            performance = "Acceptable (within Complex target)"
        else:
            performance = "Above target - requires attention"

        insights.append({
            'cdg': cdg,
            'current_median': current_median,
            'rolling_median': rolling_median,
            'trend': trend,
            'performance': performance,
            'month_values': current_period_values
        })

    # Sort by current median (worst first)
    insights.sort(key=lambda x: x['current_median'], reverse=True)

    # Create narrative text boxes
    y_pos = 1.2

    # Header with time period (LARGER TEXT, NHS BLUE)
    add_text_box(slide, f"Current Period ({period_label}) vs 12-Month Rolling Average",
                left=0.8, top=y_pos, width=11.7, height=0.5,
                font_size=16, bold=True, color_rgb=NHS_COLORS_RGB['nhs_blue'])
    y_pos += 0.6

    # Show top 5 CDGs
    for idx, insight in enumerate(insights[:5]):
        # Determine color based on performance
        if insight['current_median'] <= 40:
            box_color = NHS_COLORS_RGB['nhs_green']
            bg_color = RGBColor(240, 252, 240)  # Very light green tint
        elif insight['current_median'] <= 65:
            box_color = NHS_COLORS_RGB['nhs_warm_yellow']
            bg_color = RGBColor(255, 253, 240)  # Very light yellow tint
        else:
            box_color = NHS_COLORS_RGB['nhs_dark_red']  # Consistent with other slides
            bg_color = RGBColor(255, 245, 245)  # Very light red tint

        # Create colored box (FIT TO PAGE)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos),
            Inches(11.7), Inches(0.95)  # Sized to fit 5 boxes on page
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = RGBColor(*box_color)
        box.line.width = Pt(3)

        text_frame = box.text_frame
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_top = Inches(0.08)
        text_frame.word_wrap = True

        # CDG name (bold, LARGER TEXT)
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{idx+1}. {insight['cdg'][:50]}\n"
        run.font.size = Pt(14)  # Readable but fits
        run.font.bold = True
        run.font.color.rgb = RGBColor(*box_color)

        # Current period median (LARGER TEXT)
        run = p.add_run()
        run.text = f"Current period median: {insight['current_median']:.1f} days • 12-month rolling median: {insight['rolling_median']:.1f} days\n"
        run.font.size = Pt(12)  # Readable but fits
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_black'])

        # Trend comparison (LARGER TEXT)
        run = p.add_run()
        if insight['trend'] > 0:
            trend_text = f"⚠ Worsening: {abs(insight['trend']):.1f} days ABOVE 12-month average"
        elif insight['trend'] < 0:
            trend_text = f"✓ Improving: {abs(insight['trend']):.1f} days BELOW 12-month average"
        else:
            trend_text = "Stable: matching 12-month average"
        run.text = f"{trend_text} • {insight['performance']}"
        run.font.size = Pt(10)  # Readable but fits
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_mid_grey'])

        y_pos += 1.05  # Box height + gap to fit on page

    # Summary text at bottom (LARGER TEXT)
    if len(insights) > 5:
        add_text_box(slide, f"Showing top 5 of {len(insights)} CDGs by current median. Targets: Basic=25d, Regular=40d, Complex=65d",
                    left=0.8, top=6.5, width=11.7, height=0.5, font_size=11,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.LEFT)
    else:
        add_text_box(slide, "Targets: Basic=25 days, Regular=40 days, Complex=65 days",
                    left=0.8, top=6.5, width=11.7, height=0.5, font_size=11,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.LEFT)

def create_slide_12month_specialty_trends(prs, hpeg_name, trends_12month, hpeg_specialties, metrics, temp_dir):
    """NEW Slide: Current Period vs 12-Month Rolling Average - Specialties."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Resolution Time Trends by Specialty")

    # Get current period months from dataframe
    df_current = metrics.get('dataframe_current')
    if df_current is None or len(df_current) == 0:
        add_text_box(slide, "No data available",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    # Get the 3 months in the current reporting period
    if 'Completed Date' in df_current.columns:
        df_temp = df_current.copy()
        df_temp['Completed Month'] = pd.to_datetime(df_temp['Completed Date'], errors='coerce').dt.to_period('M')
        current_months = sorted(df_temp['Completed Month'].dropna().unique())[-3:]  # Last 3 months
    else:
        add_text_box(slide, "No completion date data available",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    if len(current_months) == 0:
        add_text_box(slide, "No completion data for current period",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    # Format month labels
    month_labels = [m.strftime('%b %y') for m in current_months]
    period_label = f"{month_labels[0]}-{month_labels[-1]}"

    # Analyze trend data
    specialty_trends = trends_12month.get('resolution_by_specialty_monthly', {})
    specialty_trends_filtered = {spec: months for spec, months in specialty_trends.items() if spec in hpeg_specialties}

    if len(specialty_trends_filtered) == 0:
        add_text_box(slide, "No specialty trend data available for this HPEG",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)
        return

    insights = []
    for specialty, monthly_data in specialty_trends_filtered.items():
        # Get last 12 months for rolling average
        cutoff_date = pd.Period('2024-12', 'M')
        filtered_months = {m: v for m, v in monthly_data.items() if pd.Period(m) >= cutoff_date}
        sorted_months = sorted(filtered_months.keys())[-12:]  # Last 12 months
        twelve_month_values = [filtered_months[m] for m in sorted_months]

        # Calculate 12-month rolling median
        rolling_median = np.median(twelve_month_values) if len(twelve_month_values) > 0 else 0

        # Get current 3-month values
        current_period_values = [monthly_data.get(m, None) for m in current_months]
        current_period_values = [v for v in current_period_values if v is not None and v > 0]

        if len(current_period_values) == 0:
            continue

        # Calculate current period median
        current_median = np.median(current_period_values)

        # Calculate difference
        trend = current_median - rolling_median

        # Determine performance
        if current_median <= 25:
            performance = "Excellent (within Basic target)"
        elif current_median <= 40:
            performance = "Good (within Regular target)"
        elif current_median <= 65:
            performance = "Acceptable (within Complex target)"
        else:
            performance = "Above target - requires attention"

        insights.append({
            'specialty': specialty,
            'current_median': current_median,
            'rolling_median': rolling_median,
            'trend': trend,
            'performance': performance
        })

    # Sort by current median (worst first)
    insights.sort(key=lambda x: x['current_median'], reverse=True)

    # Create narrative text boxes
    y_pos = 1.2

    # Header with time period (LARGER TEXT, NHS BLUE)
    add_text_box(slide, f"Current Period ({period_label}) vs 12-Month Rolling Average",
                left=0.8, top=y_pos, width=11.7, height=0.5,
                font_size=16, bold=True, color_rgb=NHS_COLORS_RGB['nhs_blue'])
    y_pos += 0.6

    # Show top 5 specialties
    for idx, insight in enumerate(insights[:5]):
        # Determine color based on performance
        if insight['current_median'] <= 40:
            box_color = NHS_COLORS_RGB['nhs_green']
            bg_color = RGBColor(240, 252, 240)  # Very light green tint
        elif insight['current_median'] <= 65:
            box_color = NHS_COLORS_RGB['nhs_warm_yellow']
            bg_color = RGBColor(255, 253, 240)  # Very light yellow tint
        else:
            box_color = NHS_COLORS_RGB['nhs_dark_red']  # Consistent with other slides
            bg_color = RGBColor(255, 245, 245)  # Very light red tint

        # Create colored box (FIT TO PAGE)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos),
            Inches(11.7), Inches(0.95)  # Sized to fit 5 boxes on page
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = RGBColor(*box_color)
        box.line.width = Pt(3)

        text_frame = box.text_frame
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_top = Inches(0.08)
        text_frame.word_wrap = True

        # Specialty name (bold, LARGER TEXT)
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{idx+1}. {insight['specialty'][:50]}\n"
        run.font.size = Pt(14)  # Readable but fits
        run.font.bold = True
        run.font.color.rgb = RGBColor(*box_color)

        # Current period median (LARGER TEXT)
        run = p.add_run()
        run.text = f"Current period median: {insight['current_median']:.1f} days • 12-month rolling median: {insight['rolling_median']:.1f} days\n"
        run.font.size = Pt(12)  # Readable but fits
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_black'])

        # Trend comparison (LARGER TEXT)
        run = p.add_run()
        if insight['trend'] > 0:
            trend_text = f"⚠ Worsening: {abs(insight['trend']):.1f} days ABOVE 12-month average"
        elif insight['trend'] < 0:
            trend_text = f"✓ Improving: {abs(insight['trend']):.1f} days BELOW 12-month average"
        else:
            trend_text = "Stable: matching 12-month average"
        run.text = f"{trend_text} • {insight['performance']}"
        run.font.size = Pt(10)  # Readable but fits
        run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_mid_grey'])

        y_pos += 1.05  # Box height + gap to fit on page

    # Summary text at bottom (LARGER TEXT)
    if len(insights) > 5:
        add_text_box(slide, f"Showing top 5 of {len(insights)} specialties by current median. Targets: Basic=25d, Regular=40d, Complex=65d",
                    left=0.8, top=6.5, width=11.7, height=0.5, font_size=11,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.LEFT)
    else:
        add_text_box(slide, "Targets: Basic=25 days, Regular=40 days, Complex=65 days",
                    left=0.8, top=6.5, width=11.7, height=0.5, font_size=11,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.LEFT)

def create_slide_resolution_complexity_distribution(prs, hpeg_name, metrics, temp_dir):
    """NEW Slide: Resolution Time Distribution by Complexity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Resolution Time by Complexity Level")

    # Get current period months from dataframe for time period label
    df_current = metrics.get('dataframe_current')
    period_label = "Current Period"
    if df_current is not None and len(df_current) > 0 and 'Completed Date' in df_current.columns:
        df_temp = df_current.copy()
        df_temp['Completed Month'] = pd.to_datetime(df_temp['Completed Date'], errors='coerce').dt.to_period('M')
        current_months = sorted(df_temp['Completed Month'].dropna().unique())[-3:]  # Last 3 months
        if len(current_months) > 0:
            month_labels = [m.strftime('%b %y') for m in current_months]
            period_label = f"{month_labels[0]}-{month_labels[-1]}"

    # Add summary boxes showing mean/median by complexity
    y_pos = 1.2
    resolution_by_comp = metrics.get('resolution_by_complexity', {})

    if len(resolution_by_comp) > 0:
        # Header with time period (LARGER TEXT, NHS BLUE)
        add_text_box(slide, f"Resolution Times & Performance Metrics for {period_label} (Business Days)", left=0.8, top=y_pos,
                    width=11.7, height=0.4, font_size=16, bold=True,
                    color_rgb=NHS_COLORS_RGB['nhs_blue'])
        y_pos += 0.6

        # Table-like layout - LARGER boxes using more space
        complexities = ['Basic', 'Regular', 'Complex']
        targets = [25, 40, 65]
        colors = [NHS_COLORS_RGB['nhs_green'], NHS_COLORS_RGB['nhs_warm_yellow'],
                  NHS_COLORS_RGB['nhs_dark_red']]  # Consistent with other slides

        box_width = 3.7
        box_height = 1.5
        spacing = 0.25
        start_x = 0.9

        for idx, (complexity, target, color) in enumerate(zip(complexities, targets, colors)):
            if complexity in resolution_by_comp:
                data = resolution_by_comp[complexity]
                mean_val = data.get('mean', 0)
                median_val = data.get('median', 0)
                count = data.get('count', 0)

                x_pos = start_x + (idx * (box_width + spacing))

                # Create colored box with subtle tint
                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(box_width), Inches(box_height)
                )
                box.fill.solid()
                # Light tinted backgrounds matching slides 8 & 9
                if idx == 0:
                    box.fill.fore_color.rgb = RGBColor(240, 252, 240)  # Light green
                elif idx == 1:
                    box.fill.fore_color.rgb = RGBColor(255, 253, 240)  # Light yellow
                else:
                    box.fill.fore_color.rgb = RGBColor(255, 245, 245)  # Light red

                box.line.color.rgb = RGBColor(*color)
                box.line.width = Pt(3)

                # Add text to box
                text_frame = box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_top = Inches(0.1)

                # Title (LARGER TEXT)
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = f"{complexity}\n"
                run.font.size = Pt(15)  # Increased from 12 to 15
                run.font.bold = True
                run.font.color.rgb = RGBColor(*color)

                # Mean (LARGER TEXT)
                p.add_run()
                run = p.add_run()
                run.text = f"Mean: {mean_val:.1f} days\n"
                run.font.size = Pt(13)  # Increased from 10 to 13
                run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_black'])

                # Median (LARGER TEXT)
                run = p.add_run()
                run.text = f"Median: {median_val:.1f} days\n"
                run.font.size = Pt(13)  # Increased from 10 to 13
                run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_black'])

                # Count (LARGER TEXT)
                run = p.add_run()
                run.text = f"Cases: {count}\n"
                run.font.size = Pt(11)  # Increased from 9 to 11
                run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_mid_grey'])

                # Target (LARGER TEXT)
                run = p.add_run()
                run.text = f"Target: {target} days"
                run.font.size = Pt(11)  # Increased from 9 to 11
                run.font.bold = True
                run.font.color.rgb = RGBColor(*color)

        # Add performance summary chart below boxes
        y_bottom = y_pos + box_height + 0.4
        targets_met = metrics.get('targets_met_by_complexity', {})

        # Performance Summary Section (LARGER TEXT, NHS BLUE)
        add_text_box(slide, "Performance Summary", left=0.8, top=y_bottom,
                    width=11.7, height=0.4, font_size=14, bold=True,
                    color_rgb=NHS_COLORS_RGB['nhs_blue'])
        y_bottom += 0.5

        # Create performance bars for each complexity
        for idx, complexity in enumerate(complexities):
            if complexity in resolution_by_comp:
                data = resolution_by_comp[complexity]
                mean_val = data.get('mean', 0)
                target = targets[idx]
                compliance_pct = targets_met.get(complexity, 0)
                count = data.get('count', 0)

                # Performance bar box with subtle background tint
                perf_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.9), Inches(y_bottom),
                    Inches(11.5), Inches(0.65)
                )
                perf_box.fill.solid()
                # Subtle tinted backgrounds
                if idx == 0:
                    perf_box.fill.fore_color.rgb = RGBColor(245, 253, 245)  # Very light green
                elif idx == 1:
                    perf_box.fill.fore_color.rgb = RGBColor(255, 254, 245)  # Very light yellow
                else:
                    perf_box.fill.fore_color.rgb = RGBColor(255, 248, 248)  # Very light red
                perf_box.line.color.rgb = RGBColor(*colors[idx])
                perf_box.line.width = Pt(2)

                text_frame = perf_box.text_frame
                text_frame.margin_left = Inches(0.15)
                text_frame.margin_top = Inches(0.1)
                text_frame.word_wrap = True

                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                # Complexity name (LARGER TEXT)
                run = p.add_run()
                run.text = f"{complexity}: "
                run.font.size = Pt(13)  # Increased from 11 to 13
                run.font.bold = True
                run.font.color.rgb = RGBColor(*colors[idx])

                # Performance details (LARGER TEXT)
                run = p.add_run()
                performance_status = "✓ Meeting target" if mean_val <= target else f"⚠ {mean_val - target:.1f} days over target"
                run.text = f"{count} cases • Avg: {mean_val:.1f} days • Target: {target} days • {performance_status} • {compliance_pct:.0f}% compliant"
                run.font.size = Pt(12)  # Increased from 10 to 12
                run.font.color.rgb = RGBColor(*NHS_COLORS_RGB['nhs_black'])

                y_bottom += 0.75

        # Add overall summary at bottom (LARGER TEXT)
        total_cases = sum(resolution_by_comp[c].get('count', 0) for c in complexities if c in resolution_by_comp)
        overall_compliance = sum(targets_met.get(c, 0) * resolution_by_comp[c].get('count', 0)
                               for c in complexities if c in resolution_by_comp)
        overall_compliance = overall_compliance / total_cases if total_cases > 0 else 0

        add_text_box(slide, f"Overall: {total_cases} total cases • {overall_compliance:.1f}% overall compliance",
                    left=0.8, top=6.3, width=11.7, height=0.4, font_size=12, bold=True,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.CENTER)

        # Add targets footer text
        add_text_box(slide, "Targets: Basic=25d, Regular=40d, Complex=65d",
                    left=0.8, top=6.8, width=11.7, height=0.3, font_size=10,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.CENTER)
    else:
        add_text_box(slide, "No resolution time data available for current period",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)

def create_slide_seasonal_patterns(prs, hpeg_name, trends_12month, metrics, temp_dir):
    """NEW Slide: Seasonal Volume Patterns (HPEG-specific)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Seasonal Patterns & Volume Trends")

    # Calculate HPEG-specific volume from dataframe instead of using trust-wide data
    df_current = metrics.get('dataframe_current')
    volume_monthly = {}

    if df_current is not None and len(df_current) > 0 and 'First Received' in df_current.columns:
        # Calculate monthly received complaints for THIS HPEG ONLY
        df_temp = df_current.copy()
        df_temp['Received Month'] = pd.to_datetime(df_temp['First Received']).dt.to_period('M')
        volume_monthly = df_temp.groupby('Received Month').size().to_dict()

    if len(volume_monthly) > 0:
        # Create chart
        seasonal_fig = create_seasonal_volume_chart(volume_monthly)
        seasonal_path = temp_dir / f"seasonal_{hpeg_name}.png"
        seasonal_fig.savefig(seasonal_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
        plt.close(seasonal_fig)

        # Add chart to slide
        add_image_to_slide(slide, seasonal_path, left=0.6, top=1.5, width=11.8)

        # Add interpretation text
        add_text_box(slide, "Analysis: Trend line shows overall trajectory. Look for seasonal peaks (e.g., winter pressures) and dips.",
                    left=0.8, top=6.5, width=11.7, height=0.5, font_size=9,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.LEFT)
    else:
        add_text_box(slide, "No volume trend data available",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)

def create_slide_subject_deepdive(prs, hpeg_name, metrics, temp_dir):
    """NEW Slide: Subject Deep-Dive Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, hpeg_name, "Subject Deep-Dive: Top 3 Complaint Themes")

    df_current = metrics.get('dataframe_current')

    if df_current is not None and len(df_current) > 0 and 'Subjects' in df_current.columns:
        # Create chart
        deepdive_fig = create_subject_deepdive_chart(df_current, top_n=3)
        deepdive_path = temp_dir / f"subject_deepdive_{hpeg_name}.png"
        deepdive_fig.savefig(deepdive_path, dpi=120, bbox_inches='tight', facecolor='white', pad_inches=0.1)
        plt.close(deepdive_fig)

        # Add chart to slide
        add_image_to_slide(slide, deepdive_path, left=0.6, top=1.5, width=11.8)

        # Add interpretation
        add_text_box(slide, "Each panel shows the top 5 CDGs contributing to that subject. Use this to identify which teams need targeted interventions.",
                    left=0.8, top=6.5, width=11.7, height=0.5, font_size=9,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'], alignment=PP_ALIGN.LEFT)
    else:
        add_text_box(slide, "No subject data available for current period",
                    left=2, top=3, width=9.33, height=2, font_size=12,
                    color_rgb=NHS_COLORS_RGB['nhs_mid_grey'],
                    alignment=PP_ALIGN.CENTER)

# ============================================================================
# MAIN REPORT GENERATION
# ============================================================================

def generate_hpeg_report(hpeg_name, data, temp_dir):
    """Generate complete PowerPoint report for one HPEG."""
    print(f"\n  Generating report for {hpeg_name}...")

    # Get data
    metrics = data['hpeg_metrics'][hpeg_name]

    # Skip if no complaints
    if metrics['total_current'] == 0:
        print(f"    ⚠ Skipping {hpeg_name} - no complaints in current period")
        return None

    topic_data = data['topic_models']['trust_wide']
    hpeg_dist = topic_data['hpeg_distributions'][hpeg_name]
    # Get demographic findings (if any exist and have HPEG association)
    demographic_findings = []
    if len(data['demographic_findings']) > 0:
        if isinstance(data['demographic_findings'][0], dict) and 'hpeg' in data['demographic_findings'][0]:
            demographic_findings = [f for f in data['demographic_findings'] if f.get('hpeg') == hpeg_name]
        else:
            demographic_findings = data['demographic_findings']

    # Create presentation
    prs = create_presentation()

    # Get 12-month trends data from data
    trends_12month = data.get('trends_12month', {})

    # Extract HPEG-specific CDGs and specialties from current dataframe
    df_current = metrics.get('dataframe_current')
    hpeg_cdgs = []
    hpeg_specialties = []
    if df_current is not None and len(df_current) > 0:
        if 'CDG' in df_current.columns:
            hpeg_cdgs = df_current['CDG'].dropna().unique().tolist()
        if 'Specialty' in df_current.columns:
            hpeg_specialties = df_current['Specialty'].dropna().unique().tolist()

    # Generate all slides
    create_slide_1_executive_dashboard(prs, hpeg_name, metrics, temp_dir)
    create_slide_2_locations(prs, hpeg_name, metrics, temp_dir)  # Slide 2: Locations
    create_slide_3_specialties(prs, hpeg_name, metrics, temp_dir)  # Slide 3: Specialties
    create_slide_cdg_breakdown(prs, hpeg_name, metrics, temp_dir)  # Slide 4: CDG Breakdown
    create_slide_3_rolling_subjects(prs, hpeg_name, metrics, temp_dir)  # Slide 5: Subject Trends
    create_slide_narrative_insights(prs, hpeg_name, metrics)  # Slide 6: Narrative Insights

    # NEW: Topic Intelligence & Priority Themes (Slide 7)
    topic_priorities = data.get('topic_analysis', {}).get(hpeg_name, [])
    create_slide_4_topic_intelligence(prs, hpeg_name, topic_priorities, temp_dir)  # Slide 7: Topic Intelligence

    # NEW ANALYSIS SLIDES
    create_slide_response_time_analysis(prs, hpeg_name, data['hpeg_metrics'], temp_dir)  # Slide 8: Response Time (ALL HPEGs comparison)
    create_slide_12month_cdg_trends(prs, hpeg_name, trends_12month, hpeg_cdgs, metrics, temp_dir)  # Slide 9: CDG 12-month trends
    create_slide_12month_specialty_trends(prs, hpeg_name, trends_12month, hpeg_specialties, metrics, temp_dir)  # Slide 10: Specialty 12-month trends
    create_slide_resolution_complexity_distribution(prs, hpeg_name, metrics, temp_dir)  # Slide 11: Complexity Resolution
    create_slide_seasonal_patterns(prs, hpeg_name, trends_12month, metrics, temp_dir)  # Slide 12: Seasonal Patterns (HPEG-specific)
    create_slide_subject_deepdive(prs, hpeg_name, metrics, temp_dir)  # Slide 13: Subject Deep-Dive

    # EXISTING SLIDES (renumbered)
    create_slide_5_performance_metrics(prs, hpeg_name, metrics, temp_dir)  # Slide 14: Complexity Donut
    create_slide_6_risk_dashboard(prs, hpeg_name, metrics)  # Slide 15: Risk Dashboard
    create_slide_7_demographics(prs, hpeg_name, metrics, demographic_findings)  # Slide 16: Demographics (conditional)
    create_slide_8_oldest_cases(prs, hpeg_name)  # Slide 17: Oldest Cases (blank)
    create_slide_9_actions_status(prs, hpeg_name)  # Slide 18: Actions (blank)
    create_slide_10_current_performance(prs, hpeg_name)  # Slide 19: Performance (blank)

    # Save
    period_label = data['periods']['label_current'].replace(' ', '-')
    filename = f"{hpeg_name.replace(' ', '_')}_Complaints_{period_label}.pptx"
    output_path = Path(OUTPUT_FOLDER) / filename
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs.save(str(output_path))
    print(f"    ✓ Saved: {filename}")

    return output_path

def main():
    """Main execution."""
    print("\n" + "="*70)
    print(" HPEG REPORT GENERATOR")
    print(" NHS Complaints Reporting - Script 2 of 2")
    print("="*70)

    # Load processed data
    data_path = Path(PROCESSED_DATA_PATH)
    if not data_path.exists():
        print(f"\n✗ ERROR: Processed data not found at: {PROCESSED_DATA_PATH}")
        print("Please run hpeg_data_processor.py first.")
        sys.exit(1)

    print(f"\nLoading processed data from: {data_path.name}")
    with open(data_path, 'rb') as f:
        data = pickle.load(f)

    print(f"  ✓ Data loaded successfully")
    print(f"  ✓ Period: {data['periods']['label_current']}")
    print(f"  ✓ Total complaints: {data['metadata']['total_complaints']:,}")

    # Create temporary directory for charts
    temp_dir = Path(tempfile.mkdtemp())
    print(f"  ✓ Temporary directory: {temp_dir}")

    # Generate reports for all 6 HPEGs
    print(f"\n{'='*70}")
    print("GENERATING POWERPOINT REPORTS")
    print(f"{'='*70}")

    hpegs = ['BHH Exec Team', 'QEH Exec Team', 'GHH Exec Team',
            'SH Exec Team', 'W&C Exec Team', 'CSS Exec Team']

    generated_files = []
    for hpeg in hpegs:
        try:
            output_path = generate_hpeg_report(hpeg, data, temp_dir)
            if output_path is not None:
                generated_files.append(output_path)
        except Exception as e:
            print(f"    ✗ Error generating {hpeg}: {e}")
            import traceback
            traceback.print_exc()

    # Cleanup temp directory
    import shutil
    shutil.rmtree(temp_dir)

    # Summary
    print(f"\n{'='*70}")
    print(" REPORT GENERATION COMPLETE")
    print(f"{'='*70}")
    print(f"\nGenerated {len(generated_files)} reports:")
    for path in generated_files:
        print(f"  ✓ {path.name}")

    print(f"\nAll reports saved to: {OUTPUT_FOLDER}")
    print("="*70 + "\n")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n✗ Process cancelled by user")
        sys.exit(0)
    except Exception as e:
        print(f"\n✗ ERROR: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
